import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Core Color Palette (Sleek Dark Theme)
    BG_COLOR = RGBColor(15, 23, 42)      # Slate 900
    CARD_BG = RGBColor(30, 41, 59)       # Slate 800
    TEXT_LIGHT = RGBColor(241, 245, 249)  # Slate 100
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400
    ACCENT_SKY = RGBColor(14, 165, 233)  # Sky 500
    ACCENT_ROSE = RGBColor(244, 63, 94)   # Rose 500
    ACCENT_TEAL = RGBColor(20, 184, 166)  # Teal 500

    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Top decorative accent line
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = ACCENT_SKY
        top_bar.line.fill.background()

    def add_title(slide, title_text):
        tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(1.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Segoe UI"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT

    def create_card(slide, left, top, width, height, title, content_list, title_color=ACCENT_SKY):
        # Background card shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG
        
        # Text frame
        tx_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        
        # Card Title
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.space_after = Pt(12)
        
        # Card Content
        for idx, item in enumerate(content_list):
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.name = "Segoe UI"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_LIGHT
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: Title & Executive Summary
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6]  # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1)
    
    # Large Cover Title
    cover_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.833), Inches(2.2))
    tf = cover_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    
    p1 = tf.paragraphs[0]
    p1.text = "FictiPay Customer Churn Prediction System"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_after = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = "Proactive Profit Protection via High-Scale Distributed ML & Explainable AI (XAI)"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_SKY
    
    # Executive Summary Cards
    create_card(
        slide1, 
        left=Inches(0.75), top=Inches(3.8), width=Inches(3.6), height=Inches(2.8),
        title="The Retention Problem",
        title_color=ACCENT_ROSE,
        content_list=[
            "Acquisition costs 5x to 25x more than retaining a wallet user.",
            "Silent churn causes massive revenue erosion before detection.",
            "Proactive mitigation preserves customer lifetime value (CLTV)."
        ]
    )
    
    create_card(
        slide1, 
        left=Inches(4.8), top=Inches(3.8), width=Inches(3.6), height=Inches(2.8),
        title="The Predictive Objective",
        title_color=ACCENT_TEAL,
        content_list=[
            "Identify high-risk accounts 30 days before silent churn occurs.",
            "Translate model explanations into targeted local interventions.",
            "Orchestrate an automated win-back rule engine for marketing."
        ]
    )
    
    create_card(
        slide1, 
        left=Inches(8.85), top=Inches(3.8), width=Inches(3.7), height=Inches(2.8),
        title="Key System Outcomes",
        title_color=ACCENT_SKY,
        content_list=[
            "Robust 10-fold CV ensemble achieves a peak 0.983 AUC score.",
            "SHAP TreeExplainer ensures leakage-free customer attribution.",
            "Threshold optimization balances precision and intervention cost."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 2: Scalable Data Architecture
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2)
    add_title(slide2, "Scalable Data Architecture: Handling Massive Scale")
    
    create_card(
        slide2, 
        left=Inches(0.75), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Framework Strategy: Dask vs. Spark",
        title_color=ACCENT_SKY,
        content_list=[
            "Selected Dask for out-of-core, distributed computing in Python.",
            "Deep integration with the PyData stack (Scikit-Learn, XGBoost).",
            "Lighter-weight cluster orchestration compared to JVM-based Spark.",
            "Avoids expensive serialization layers, keeping debugging native.",
            "Handles multi-gigabyte datasets smoothly on single-node or clusters."
        ]
    )
    
    create_card(
        slide2, 
        left=Inches(6.98), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Partitioning & Bucketing Design",
        title_color=ACCENT_TEAL,
        content_list=[
            "Data Scales: KYC (2M rows), Transactions (200M+ rows), Day-End Balances (360M rows).",
            "Lazy Evaluation: Read Parquet datasets using strict column pruning.",
            "Month-based partitioning isolates historical windows efficiently.",
            "Bucketing on ACCOUNT_ID aligns partitions across transactions and balance tables.",
            "Minimizes expensive distributed network shuffling during massive merges."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 3: Rolling-Window Framework & Feature Engineering
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3)
    add_title(slide3, "Rolling-Window Framework & Feature Quality")
    
    create_card(
        slide3, 
        left=Inches(0.75), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Temporal Rolling-Window Framework",
        title_color=ACCENT_SKY,
        content_list=[
            "Separates historical inputs from future outcomes to prevent leakage.",
            "Observation Window: Jan 1 – Mar 31 (90 days behavioral logs).",
            "Evaluation Window: Apr 1 – Apr 30 (30 days inactivity window).",
            "Target Churn: CHURN = 1 if transactions in April == 0, else 0.",
            "Allows detection of silent churners who slowly drain their wallets.",
            "Strict cutoff at March 31 ensures zero leakage of future state."
        ]
    )
    
    create_card(
        slide3, 
        left=Inches(6.98), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Handling Data Skew & Zero-Inflation",
        title_color=ACCENT_TEAL,
        content_list=[
            "Financial aggregates (spend, counts) exhibit extreme log-normal skewness.",
            "Applied log1p transformations to regularize feature distributions.",
            "Zero-Inflation: High percentage of zero values in specific services (e.g., BillPay).",
            "Used a two-stage Hurdle Model approach: Flagging zero-state vs. modeling quantity.",
            "Adding binary zero-activity flags preserves physical meaning for tree nodes."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 4: Mitigating Class Imbalance
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide4)
    add_title(slide4, "Mitigating Class Imbalance: Focusing on the Minority")
    
    create_card(
        slide4, 
        left=Inches(0.75), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Traditional Resampling vs. Hybrid SMOTE-ENN",
        title_color=ACCENT_SKY,
        content_list=[
            "Severe class imbalance present with Churn prevalence at ~5%.",
            "Traditional SMOTE creates synthetic points but introduces noisy outliers.",
            "SMOTE-ENN oversamples the minority class and then uses Edited Nearest Neighbors to clean the decision boundary.",
            "Removes overlapping majority samples, sharpening structural boundaries.",
            "Enhances precision near complex boundaries, reducing false positives."
        ]
    )
    
    create_card(
        slide4, 
        left=Inches(6.98), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Class-Weighted Loss & Scalability",
        title_color=ACCENT_TEAL,
        content_list=[
            "For massive terabyte-scale datasets, SMOTE-ENN is too slow.",
            "Alternative: Adjust class weights directly in the loss function.",
            "Set scale_pos_weight in XGBoost/LightGBM dynamically based on imbalance ratio.",
            "Forces gradient updates to prioritize errors on the minority class.",
            "Maintains O(N) training scale while yielding robust recall profiles."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 5: Advanced Ensemble Modeling
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide5)
    add_title(slide5, "Advanced Ensemble: Dual-Engine Model Zoo")
    
    create_card(
        slide5, 
        left=Inches(0.75), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Ensemble Architecture",
        title_color=ACCENT_SKY,
        content_list=[
            "Combines tabular boosters with deep representation models.",
            "XGBoost & LightGBM: Build deep, high-efficiency decision splits.",
            "FT-Transformer: Encodes categorical/numerical inputs through self-attention layers to capture complex cross-feature interactions.",
            "Diversity in model structures stabilizes predictions against noise.",
            "Peak Out-of-Fold (OOF) validation ROC-AUC reaches 0.983."
        ]
    )
    
    create_card(
        slide5, 
        left=Inches(6.98), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Out-of-Fold Stacking & Calibration",
        title_color=ACCENT_TEAL,
        content_list=[
            "Used 10-Fold Stratified Cross-Validation to generate OOF outputs.",
            "A Logistic Regression meta-learner blends base model predictions.",
            "Rank-Average Blending converts probabilities to percentiles to avoid calibration variance.",
            "Isotonic Regression calibrates blended ranks back to true probability.",
            "Secures smooth, monotonically calibrated risk outputs for downstream tools."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 6: Explainable AI (XAI) with SHAP
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide6)
    add_title(slide6, "Explainable AI (XAI): Model Transparency via SHAP")
    
    create_card(
        slide6, 
        left=Inches(0.75), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Global Interpretability (Beeswarm)",
        title_color=ACCENT_SKY,
        content_list=[
            "Calculates exact Shapley values to attribute feature importance.",
            "Top Churn Risk Drivers: Outbound transaction recency, increasing daily balance volatility, and decaying activity trends.",
            "Top Protective Drivers: Long account tenure, usage of sticky features (e.g. BillPay), and positive monthly net-flow.",
            "Reveals non-linear trends (e.g. very high balance is protective, but volatile balance indicates flight risk)."
        ]
    )
    
    create_card(
        slide6, 
        left=Inches(6.98), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Local Interpretability (Waterfall)",
        title_color=ACCENT_TEAL,
        content_list=[
            "Explains the exact churn probability for any individual customer.",
            "Traces base value (dataset average) to prediction value via additive step-wise SHAP contributions.",
            "Example: High-risk customer's 82% churn probability driven by: monthly balance drop (+24% risk), month-to-month contract flag (+18% risk), and zero outbound transactions in March (+15% risk)."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 7: Translating Predictions to Win-Back Strategies
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide7)
    add_title(slide7, "Translating Predictions to Targeted Win-Back Strategies")
    
    create_card(
        slide7, 
        left=Inches(0.75), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Rule-Based Automated Interventions",
        title_color=ACCENT_SKY,
        content_list=[
            "Recency Gap (High Inactivity) -> Offer transaction cashback vouchers.",
            "Draining Balance (Negative Trend) -> Interest-booster balance campaign.",
            "Zero Bill Pay / Merchant Pay -> Utility discount or merchant store vouchers (establishes high-stickiness wallet habits).",
            "High Activity Decay (Sudden Drops) -> Initiate proactive outbound customer service calls to resolve support friction points."
        ]
    )
    
    create_card(
        slide7, 
        left=Inches(6.98), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Precision Campaigns & Fatigue Mitigation",
        title_color=ACCENT_TEAL,
        content_list=[
            "Prevents customer alert fatigue: Only target the top 10% highest-risk tier.",
            "Optimizes marketing spend by ignoring stable customers who don't need incentives.",
            "Omits lost-cause churners (extremely low balance, inactive for 90+ days).",
            "Creates tailored, automatic customer journeys fed directly by nightly SHAP inference."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 8: Financial ROI & Conclusion
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide8)
    add_title(slide8, "Financial ROI & Conclusion")
    
    create_card(
        slide8, 
        left=Inches(0.75), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Cost-Sensitive Decision Optimization",
        title_color=ACCENT_SKY,
        content_list=[
            "Retention campaigns are not free. Setting a decision threshold is a balance.",
            "Intervention Cost (e.g., cashback voucher) = 1 TK.",
            "Churn Loss (missed customer retention value) = 5 TK.",
            "Optimization Loss = 5 * False Negatives + 1 * False Positives.",
            "Found optimal mathematical threshold at P >= 0.32.",
            "Maximizes profit recovery relative to marketing outreach costs."
        ]
    )
    
    create_card(
        slide8, 
        left=Inches(6.98), top=Inches(1.8), width=Inches(5.6), height=Inches(4.8),
        title="Key Takeaways",
        title_color=ACCENT_TEAL,
        content_list=[
            "Dask handles scaling constraints securely, integrating standard Scikit-Learn structures.",
            "Hybrid SMOTE-ENN ensures sharp, robust validation splits.",
            "SHAP TreeExplainer ensures reliable, leakage-free attribution.",
            "ROI-optimized thresholding transforms predictions into a high-profit retention engine."
        ]
    )
    
    # Save the presentation
    output_filename = "FictiPay_Churn_System_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully saved as '{output_filename}'")

if __name__ == "__main__":
    create_presentation()
